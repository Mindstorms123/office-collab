# -*- coding: utf-8 -*-
import sys
import json
import os
import tempfile

sys.stdout.reconfigure(encoding='utf-8')
sys.stderr.reconfigure(encoding='utf-8')

from docx import Document
from docx.shared import Pt, RGBColor as DocxRGB
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment


# ─────────────────────────────────────────────────────────────────
# HILFSFUNKTIONEN
# ─────────────────────────────────────────────────────────────────

def hex_to_rgb(h):
    h = str(h).lstrip('#')
    if len(h) == 3:
        h = ''.join(c*2 for c in h)
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))


def fetch_image(query, width=1280, height=720):
    """Kostenloses Bild von Unsplash laden"""
    try:
        import requests
        if not query:
            return None
        url = f"https://source.unsplash.com/{width}x{height}/?{query.replace(' ', ',')}"
        r = requests.get(url, timeout=12, allow_redirects=True)
        if r.status_code == 200 and 'image' in r.headers.get('Content-Type', ''):
            f = tempfile.NamedTemporaryFile(delete=False, suffix='.jpg')
            f.write(r.content)
            f.close()
            return f.name
    except Exception as e:
        print(f'[IMG Fehler] {query}: {e}', file=sys.stderr)
    return None


def push_to_back(shape, slide):
    """Shape hinter alle anderen schieben"""
    sp = shape._element
    slide.shapes._spTree.remove(sp)
    slide.shapes._spTree.insert(2, sp)


def apply_rounded(shape, rounded):
    """Abgerundete Ecken via XML"""
    pg = shape._element.spPr.find(qn('a:prstGeom'))
    if pg is None:
        pg = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
    if pg is not None:
        pg.set('prst', 'roundRect')
        av = pg.find(qn('a:avLst'))
        if av is None:
            av = etree.SubElement(pg, qn('a:avLst'))
        av.clear()
        gd = etree.SubElement(av, qn('a:gd'))
        gd.set('name', 'adj')
        gd.set('fmla', f'val {int(rounded)}')


def apply_opacity(shape, opacity):
    """Transparenz auf Shape anwenden (0.0 = unsichtbar, 1.0 = opak)"""
    try:
        alpha_val = int(opacity * 100000)
        spPr = shape._element.spPr
        fmla = spPr.find('.//' + qn('a:solidFill'))
        if fmla is not None:
            clr = fmla.find(qn('a:srgbClr'))
            if clr is None:
                clr = fmla.find(qn('a:sysClr'))
            if clr is not None:
                # Bestehende alpha entfernen
                for a in clr.findall(qn('a:alpha')):
                    clr.remove(a)
                alpha_el = etree.SubElement(clr, qn('a:alpha'))
                alpha_el.set('val', str(alpha_val))
    except Exception as e:
        print(f'[Opacity Fehler] {e}', file=sys.stderr)


# ─────────────────────────────────────────────────────────────────
# DOCX RENDERER
# ─────────────────────────────────────────────────────────────────

def render_docx(data, output_path):
    doc = Document()

    # Styles anpassen
    style = doc.styles['Normal']
    style.font.size = Pt(11)

    doc.add_heading(data.get('title', 'Dokument'), level=0)

    for section in data.get('sections', []):
        heading  = section.get('heading')
        content  = section.get('content')
        table_data = section.get('table')

        if heading:
            doc.add_heading(heading, level=section.get('level', 1))

        if content:
            p = doc.add_paragraph(content)
            if section.get('bold'):
                for run in p.runs:
                    run.bold = True

        if table_data and len(table_data) > 0:
            t = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
            t.style = 'Table Grid'
            for i, row in enumerate(table_data):
                for j, val in enumerate(row):
                    cell = t.rows[i].cells[j]
                    cell.text = str(val)
                    if i == 0:
                        for para in cell.paragraphs:
                            for run in para.runs:
                                run.bold = True
            doc.add_paragraph('')

    doc.save(output_path)


# ─────────────────────────────────────────────────────────────────
# PPTX RENDERER – objekt-basiert, modern, mit Bildern
# ─────────────────────────────────────────────────────────────────

def render_pptx(data, output_path):
    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Palette aus KI-JSON übernehmen
    raw_palette = data.get('palette', {})
    defaults = {
        'bg1':     '#0a0f1e',
        'bg2':     '#141e3a',
        'accent1': '#38bdf8',
        'accent2': '#818cf8',
        'text':    '#f1f5f9',
        'sub':     '#64748b',
        'card':    '#0d1526',
    }
    palette = {}
    for k, fallback in defaults.items():
        palette[k] = hex_to_rgb(raw_palette.get(k, fallback))

    # Bild-Cache (gleiche Query nicht doppelt laden)
    img_cache = {}

    def get_img(query):
        if not query:
            return None
        if query not in img_cache:
            print(f'[Bild laden] {query}', file=sys.stderr)
            img_cache[query] = fetch_image(query)
        return img_cache[query]

    # ── Shape-Helfer ────────────────────────────────────────────

    def add_rect(slide, l, t, w, h, color, rounded=0, opacity=None, send_back=False):
        shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = color if isinstance(color, RGBColor) else hex_to_rgb(color)
        shape.line.fill.background()
        if rounded:
            apply_rounded(shape, rounded)
        if opacity is not None and opacity < 1.0:
            apply_opacity(shape, opacity)
        if send_back:
            push_to_back(shape, slide)
        return shape

    def add_text(slide, text, l, t, w, h, size=18, bold=False, italic=False,
                 color=None, align=PP_ALIGN.LEFT, wrap=True):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf  = txb.text_frame
        tf.word_wrap = wrap
        para = tf.paragraphs[0]
        para.alignment = align
        run = para.add_run()
        run.text       = str(text)
        run.font.size  = PptPt(size)
        run.font.bold  = bold
        run.font.italic = italic
        run.font.color.rgb = (color if isinstance(color, RGBColor) else hex_to_rgb(color)) if color else palette['text']
        return txb

    def add_image(slide, query, l, t, w, h, rounded=0, opacity=None, send_back=False):
        path = get_img(query)
        if not path or not os.path.exists(path):
            # Fallback: dunkelgraues Rechteck
            add_rect(slide, l, t, w, h, palette['bg2'], rounded=rounded)
            return
        try:
            pic = slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))
            if rounded:
                # Bild in abgerundeter Form via custGeom ist komplex –
                # stattdessen: Bild + transparentes Overlay mit abgerundeter Form
                pass
            if send_back:
                push_to_back(pic, slide)
        except Exception as e:
            print(f'[Bild einfügen Fehler] {e}', file=sys.stderr)
            add_rect(slide, l, t, w, h, palette['bg2'])

    def slide_number(slide, num):
        add_text(slide, str(num),
                 12.6, 7.05, 0.6, 0.35,
                 size=11, color=palette['sub'], align=PP_ALIGN.RIGHT)

    # ── Objekt-Dispatcher ───────────────────────────────────────

    def render_object(slide, obj):
        otype   = obj.get('type', '')
        l       = float(obj.get('l', 0))
        t       = float(obj.get('t', 0))
        w       = float(obj.get('w', 1))
        h       = float(obj.get('h', 1))
        rounded = int(obj.get('rounded', 0))
        opacity = obj.get('opacity', None)
        back    = obj.get('back', False)

        if otype == 'rect':
            add_rect(slide, l, t, w, h,
                     obj.get('color', '#000000'),
                     rounded=rounded,
                     opacity=opacity,
                     send_back=back)

        elif otype == 'text':
            align_map = {
                'left':   PP_ALIGN.LEFT,
                'center': PP_ALIGN.CENTER,
                'right':  PP_ALIGN.RIGHT,
            }
            add_text(slide,
                     obj.get('content', ''),
                     l, t, w, h,
                     size    = int(obj.get('size', 18)),
                     bold    = bool(obj.get('bold', False)),
                     italic  = bool(obj.get('italic', False)),
                     color   = obj.get('color', None),
                     align   = align_map.get(obj.get('align', 'left'), PP_ALIGN.LEFT))

        elif otype == 'image':
            add_image(slide,
                      obj.get('query', ''),
                      l, t, w, h,
                      rounded  = rounded,
                      opacity  = opacity,
                      send_back = back)

    # ── Folien rendern ──────────────────────────────────────────

    for idx, slide_data in enumerate(data.get('slides', [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        objects = slide_data.get('objects', [])

        # Hintergrund immer zuerst (back=True)
        for obj in objects:
            if obj.get('back', False):
                render_object(slide, obj)

        # Dann alle anderen in Reihenfolge
        for obj in objects:
            if not obj.get('back', False):
                render_object(slide, obj)

        # Notizen
        if slide_data.get('notes'):
            try:
                slide.notes_slide.notes_text_frame.text = slide_data['notes']
            except:
                pass

    # Temp-Bilder aufräumen
    for path in img_cache.values():
        if path:
            try:
                os.unlink(path)
            except:
                pass

    prs.save(output_path)


# ─────────────────────────────────────────────────────────────────
# XLSX RENDERER
# ─────────────────────────────────────────────────────────────────

def render_xlsx(data, output_path):
    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    header_fill = PatternFill("solid", fgColor="1E293B")
    header_font = Font(bold=True, color="38BDF8", size=11)
    accent_fill = PatternFill("solid", fgColor="0F172A")

    for sheet_data in data.get('sheets', []):
        ws = wb.create_sheet(title=sheet_data.get('name', 'Sheet'))

        headers = sheet_data.get('headers', [])
        if headers:
            ws.append(headers)
            for cell in ws[1]:
                cell.font      = header_font
                cell.fill      = header_fill
                cell.alignment = Alignment(horizontal='center', vertical='center')
            ws.row_dimensions[1].height = 24

        for row in sheet_data.get('rows', []):
            ws.append([
                v if isinstance(v, (int, float)) else (str(v) if v is not None else '')
                for v in row
            ])

        # Formeln eintragen
        for cell_ref, formula in sheet_data.get('formulas', {}).items():
            ws[cell_ref] = formula

        # Zebra-Streifen
        for row_idx, row in enumerate(ws.iter_rows(min_row=2), start=2):
            if row_idx % 2 == 0:
                for cell in row:
                    cell.fill = accent_fill

        # Spaltenbreite
        for col in ws.columns:
            max_len = max(
                (len(str(c.value)) for c in col if c.value is not None),
                default=10
            )
            ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 45)

    wb.save(output_path)


# ─────────────────────────────────────────────────────────────────
# EINSTIEGSPUNKT
# ─────────────────────────────────────────────────────────────────

if __name__ == '__main__':
    if len(sys.argv) < 3:
        print('Usage: render_office.py <output_path> <format>', file=sys.stderr)
        sys.exit(1)

    output_path = sys.argv[1]
    file_format = sys.argv[2]
    data        = json.loads(sys.stdin.buffer.read().decode('utf-8'))

    if file_format == 'docx':
        render_docx(data, output_path)
    elif file_format == 'pptx':
        render_pptx(data, output_path)
    elif file_format == 'xlsx':
        render_xlsx(data, output_path)
    else:
        print(f'Unbekanntes Format: {file_format}', file=sys.stderr)
        sys.exit(1)
